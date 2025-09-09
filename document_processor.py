"""
Document processing module for handling PDF, DOCX, PPTX, and TXT files.
Extracts text content and generates embeddings for the AI Coach system.
Now using ChromaDB for optimized vector storage and retrieval.
"""

import os
import asyncio
import logging
from typing import List, Dict, Tuple, Optional, Any
from datetime import datetime
import tempfile
import hashlib
import uuid

# Document processing imports
try:
    import pypdf
    from pdfplumber import PDF as PDFPlumber
    from docx import Document as DocxDocument
    from pptx import Presentation
    HAS_DOC_LIBS = True
except ImportError as e:
    HAS_DOC_LIBS = False
    logging.warning(f"Document processing libraries not available: {e}")

# OpenAI for embeddings and ChromaDB for vector storage
from openai import AsyncOpenAI
from chroma_client import get_chroma_client

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self, openai_client: AsyncOpenAI):
        self.openai_client = openai_client
        self.chroma_client = get_chroma_client()
        self.chunk_size = 200  # tokens per chunk - smaller for better retrieval
        self.chunk_overlap = 50  # overlap between chunks
        
    async def process_uploaded_file(self, file_path: str, filename: str, expert_name: str, session_id: str) -> Dict:
        """Process an uploaded file and store chunks with embeddings in ChromaDB"""
        try:
            # Determine file type
            file_ext = filename.lower().split('.')[-1]
            if file_ext not in ['pdf', 'docx', 'pptx', 'txt']:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            # Get file size
            file_size = os.path.getsize(file_path)
            
            logger.info(f"📄 Processing {file_ext.upper()} document: {filename}")
            
            # Extract text content
            content_chunks = await self._extract_content(file_path, file_ext)
            
            if not content_chunks:
                return {"success": False, "error": "No content extracted from document"}
            
            logger.info(f"📝 Extracted {len(content_chunks)} chunks from {filename}")
            
            # Store chunks in ChromaDB
            stored_chunks = 0
            for i, chunk_data in enumerate(content_chunks):
                try:
                    # Create unique chunk ID
                    chunk_id = f"{session_id}_{filename}_{i}_{hashlib.md5(chunk_data['content'].encode()).hexdigest()[:8]}"
                    
                    # Prepare metadata
                    metadata = {
                        'session_id': session_id,
                        'title': filename,
                        'doc_type': file_ext,
                        'expert_name': expert_name,
                        'chunk_index': i,
                        'upload_date': datetime.utcnow().isoformat(),
                        'file_size': file_size
                    }
                    
                    # Add page/slide specific metadata
                    if chunk_data.get('page_number'):
                        metadata['page_number'] = chunk_data['page_number']
                    if chunk_data.get('slide_number'):
                        metadata['slide_number'] = chunk_data['slide_number']
                    
                    # Log chunk content for debugging
                    logger.info(f"📦 CHUNK {i+1} content ({len(chunk_data['content'])} chars):")
                    logger.info(f"📝 CHUNK: {chunk_data['content'][:300]}{'...' if len(chunk_data['content']) > 300 else ''}")
                    
                    # Store in ChromaDB (embedding is generated automatically)
                    success = self.chroma_client.add_document_chunk(
                        chunk_id=chunk_id,
                        content=chunk_data['content'],
                        metadata=metadata
                    )
                    
                    if success:
                        stored_chunks += 1
                    
                except Exception as e:
                    logger.error(f"Error processing chunk {i}: {e}")
                    continue
            
            if stored_chunks > 0:
                logger.info(f"✅ Successfully processed {filename}: {stored_chunks} chunks stored in ChromaDB")
                return {
                    "success": True,
                    "chunks_processed": stored_chunks,
                    "total_chunks": len(content_chunks),
                    "filename": filename
                }
            else:
                return {"success": False, "error": "Failed to process any chunks"}
                
        except Exception as e:
            logger.error(f"Error processing document {filename}: {e}")
            return {"success": False, "error": str(e)}
    
    async def _extract_content(self, file_path: str, file_type: str) -> List[Dict]:
        """Extract text content from different file types"""
        if not HAS_DOC_LIBS:
            raise ImportError("Document processing libraries not installed")
        
        if file_type == 'pdf':
            return await self._extract_pdf_content(file_path)
        elif file_type == 'docx':
            return await self._extract_docx_content(file_path)
        elif file_type == 'pptx':
            return await self._extract_pptx_content(file_path)
        elif file_type == 'txt':
            return await self._extract_txt_content(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    async def _extract_pdf_content(self, file_path: str) -> List[Dict]:
        """Extract text from PDF using both pypdf and pdfplumber"""
        chunks = []
        
        try:
            # Try pdfplumber first (better for complex layouts)
            with PDFPlumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        # Log extracted text for debugging
                        logger.info(f"📄 PDF Page {page_num} extracted text ({len(text)} chars):")
                        logger.info(f"📝 TEXT: {text[:500]}{'...' if len(text) > 500 else ''}")
                        
                        # Split long pages into chunks
                        page_chunks = self._split_text_into_chunks(text)
                        for chunk_text in page_chunks:
                            chunks.append({
                                'content': chunk_text.strip(),
                                'page_number': page_num,
                                'source': 'pdfplumber'
                            })
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying pypdf: {e}")
            
            # Fallback to pypdf
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = pypdf.PdfReader(file)
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text()
                        if text and text.strip():
                            # Log extracted text for debugging
                            logger.info(f"📄 PDF Page {page_num} extracted text ({len(text)} chars) [pypdf]:")
                            logger.info(f"📝 TEXT: {text[:500]}{'...' if len(text) > 500 else ''}")
                            
                            page_chunks = self._split_text_into_chunks(text)
                            for chunk_text in page_chunks:
                                chunks.append({
                                    'content': chunk_text.strip(),
                                    'page_number': page_num,
                                    'source': 'pypdf'
                                })
            except Exception as e2:
                logger.error(f"Both PDF extraction methods failed: {e2}")
                raise
        
        return chunks
    
    async def _extract_docx_content(self, file_path: str) -> List[Dict]:
        """Extract text from DOCX files"""
        chunks = []
        
        try:
            doc = DocxDocument(file_path)
            full_text = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text.strip())
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text.append(' | '.join(row_text))
            
            # Join all text and split into chunks
            document_text = '\n\n'.join(full_text)
            if document_text.strip():
                # Log extracted text for debugging
                logger.info(f"📄 DOCX extracted text ({len(document_text)} chars):")
                logger.info(f"📝 TEXT: {document_text[:500]}{'...' if len(document_text) > 500 else ''}")
                
                text_chunks = self._split_text_into_chunks(document_text)
                for chunk_text in text_chunks:
                    chunks.append({
                        'content': chunk_text.strip(),
                        'source': 'docx'
                    })
        
        except Exception as e:
            logger.error(f"Error extracting DOCX content: {e}")
            raise
        
        return chunks
    
    async def _extract_pptx_content(self, file_path: str) -> List[Dict]:
        """Extract text from PowerPoint files"""
        chunks = []
        
        try:
            presentation = Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"Speaker Notes: {notes_text}")
                
                # Combine slide content
                if slide_text:
                    combined_text = '\n\n'.join(slide_text)
                    # Log extracted text for debugging
                    logger.info(f"📄 PPTX Slide {slide_num} extracted text ({len(combined_text)} chars):")
                    logger.info(f"📝 TEXT: {combined_text[:500]}{'...' if len(combined_text) > 500 else ''}")
                    
                    # Split slide content into chunks if it's too long
                    slide_chunks = self._split_text_into_chunks(combined_text)
                    for chunk_text in slide_chunks:
                        chunks.append({
                            'content': chunk_text.strip(),
                            'slide_number': slide_num,
                            'source': 'pptx'
                        })
        
        except Exception as e:
            logger.error(f"Error extracting PPTX content: {e}")
            raise
        
        return chunks
    
    async def _extract_txt_content(self, file_path: str) -> List[Dict]:
        """Extract text from TXT files"""
        chunks = []
        
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError("Could not decode text file with any common encoding")
            
            if text.strip():
                # Log extracted text for debugging
                logger.info(f"📄 TXT extracted text ({len(text)} chars):")
                logger.info(f"📝 TEXT: {text[:500]}{'...' if len(text) > 500 else ''}")
                
                text_chunks = self._split_text_into_chunks(text)
                for chunk_text in text_chunks:
                    chunks.append({
                        'content': chunk_text.strip(),
                        'source': 'txt'
                    })
        
        except Exception as e:
            logger.error(f"Error extracting TXT content: {e}")
            raise
        
        return chunks
    
    def _split_text_into_chunks(self, text: str) -> List[str]:
        """Split text into smaller chunks for embedding with better logic"""
        
        # First, try to split by double newlines (paragraphs/sections)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # Rough token estimation (1 token ≈ 4 characters)
            paragraph_tokens = len(paragraph) // 4
            current_tokens = len(current_chunk) // 4
            
            # If adding this paragraph would exceed our limit, save current and start new
            if current_tokens + paragraph_tokens > self.chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = paragraph
            else:
                if current_chunk:
                    current_chunk += "\n\n" + paragraph
                else:
                    current_chunk = paragraph
        
        # Add the last chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # If no chunks were created (very long single paragraph), fall back to sentence splitting
        if not chunks:
            sentences = text.split('. ')
            current_chunk = ""
            
            for sentence in sentences:
                estimated_tokens = len(current_chunk + sentence) // 4
                
                if estimated_tokens > self.chunk_size and current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = sentence + '. '
                else:
                    current_chunk += sentence + '. '
            
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
        
        # Ensure we have at least one chunk and handle very large single chunks
        if not chunks and text.strip():
            # For very large text, split into smaller pieces
            words = text.split()
            current_chunk = ""
            
            for word in words:
                test_chunk = current_chunk + " " + word if current_chunk else word
                if len(test_chunk) // 4 > self.chunk_size and current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = word
                else:
                    current_chunk = test_chunk
            
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
        
        return chunks if chunks else [text.strip()]
    

    def get_session_documents(self, session_id: str) -> Dict[str, Any]:
        """Get document statistics for a session"""
        try:
            return self.chroma_client.get_document_stats(session_id)
        except Exception as e:
            logger.error(f"Error getting session documents: {e}")
            return {"total_chunks": 0, "documents": []}

    def delete_session_documents(self, session_id: str) -> bool:
        """Delete all documents for a session"""
        try:
            return self.chroma_client.delete_session_documents(session_id)
        except Exception as e:
            logger.error(f"Error deleting session documents: {e}")
            return False

# Global document processor instance
document_processor = None

def get_document_processor(openai_client: AsyncOpenAI) -> DocumentProcessor:
    """Get or create document processor instance"""
    global document_processor
    if document_processor is None:
        document_processor = DocumentProcessor(openai_client)
    return document_processor